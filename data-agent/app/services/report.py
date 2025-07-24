from __future__ import annotations

from io import BytesIO
from typing import Optional

import pandas as pd
from reportlab.lib.pagesizes import letter
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches

from app.core.analysis import basic_summary, basic_insights
from app.core.charts import hist_plot


def create_pdf_report(
    df: pd.DataFrame,
    title: str = "Data Report",
    logo_path: Optional[str] = None,
) -> BytesIO:
    """Return a PDF report containing summary, insights and a sample chart."""
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter

    y = height - 40
    c.setFont("Helvetica-Bold", 16)
    c.drawCentredString(width / 2, y, title)
    y -= 40

    if logo_path:
        try:
            img = ImageReader(logo_path)
            iw, ih = img.getSize()
            aspect = ih / float(iw)
            c.drawImage(img, 40, y - 100, width=100, height=100 * aspect)
            y -= 110
        except Exception:
            pass

    c.setFont("Helvetica", 12)
    c.drawString(40, y, "Summary:")
    y -= 20
    summary = basic_summary(df)
    for k, v in summary.items():
        c.drawString(60, y, f"{k}: {v}")
        y -= 15

    y -= 10
    c.drawString(40, y, "Insights:")
    y -= 20
    insights = basic_insights(df)
    for k, v in insights.items():
        c.drawString(60, y, f"{k}: {v}")
        y -= 15

    y -= 20
    num_cols = df.select_dtypes("number").columns
    if len(num_cols) > 0:
        png = hist_plot(df, [num_cols[0]])
        img = ImageReader(png)
        iw, ih = img.getSize()
        max_width = width - 80
        scale = max_width / iw
        c.drawImage(
            img,
            40,
            y - ih * scale,
            width=iw * scale,
            height=ih * scale,
        )
        y -= ih * scale + 20

    c.showPage()
    c.save()
    buf.seek(0)
    return buf


def create_pptx_report(
    df: pd.DataFrame,
    title: str = "Data Report",
) -> BytesIO:
    """Return a PPTX report containing summary, insights and a sample chart."""
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title

    bullet = prs.slides.add_slide(prs.slide_layouts[1])
    bullet.shapes.title.text = "Summary"
    tf = bullet.shapes.placeholders[1].text_frame
    summary = basic_summary(df)
    for k, v in summary.items():
        p = tf.add_paragraph()
        p.text = f"{k}: {v}"

    p = tf.add_paragraph()
    p.text = "Insights:"
    insights = basic_insights(df)
    for k, v in insights.items():
        item = tf.add_paragraph()
        item.text = f"{k}: {v}"

    num_cols = df.select_dtypes("number").columns
    if len(num_cols) > 0:
        png = hist_plot(df, [num_cols[0]])
        chart = prs.slides.add_slide(prs.slide_layouts[5])
        chart.shapes.title.text = f"{num_cols[0]} Distribution"
        chart.shapes.add_picture(png, Inches(1), Inches(2), width=Inches(8))

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
