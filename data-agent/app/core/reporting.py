from __future__ import annotations

from typing import Any, Dict, List
from io import BytesIO
import zipfile

from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches


def history_to_pdf(history: List[Dict[str, Any]]) -> BytesIO:
    buf = BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    width, height = letter
    for item in history:
        y = height - 40
        prompt = item.get("prompt", "")
        if prompt:
            c.setFont("Helvetica-Bold", 12)
            c.drawString(40, y, f"Prompt: {prompt}")
            y -= 20
        c.setFont("Courier", 10)
        for line in item.get("code", "").splitlines():
            c.drawString(40, y, line)
            y -= 12
            if y < 100:
                c.showPage()
                y = height - 40
        for png in item.get("pngs", []):
            c.drawImage(BytesIO(png), 40, y - 200, width=400, height=200)
            y -= 210
            if y < 100:
                c.showPage()
                y = height - 40
        c.showPage()
    c.save()
    buf.seek(0)
    return buf


def history_to_pptx(history: List[Dict[str, Any]]) -> BytesIO:
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for item in history:
        slide = prs.slides.add_slide(blank)
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(9)
        height = Inches(1)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text = f"Prompt: {item.get('prompt', '')}"
        cb = slide.shapes.add_textbox(left, Inches(1.2), width, Inches(1.5))
        cb.text = item.get("code", "")
        img_top = Inches(3)
        for png in item.get("pngs", []):
            slide.shapes.add_picture(BytesIO(png), left, img_top, width=Inches(6))
            img_top += Inches(3)
    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def zip_outputs(history: List[Dict[str, Any]]) -> BytesIO:
    buf = BytesIO()
    with zipfile.ZipFile(buf, "w") as z:
        for idx, item in enumerate(history):
            for i, csv in enumerate(item.get("dfs", [])):
                z.writestr(f"table_{idx}_{i}.csv", csv)
            for i, png in enumerate(item.get("pngs", [])):
                z.writestr(f"image_{idx}_{i}.png", png)
    buf.seek(0)
    return buf
